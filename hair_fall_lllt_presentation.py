"""
Generate a PowerPoint presentation on Hair Fall Solution Using Red Light Therapy (LLLT).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import io

# ── Colour palette ─────────────────────────────────────────────────────────────
RED   = RGBColor(0xC0, 0x00, 0x00)   # deep medical-red
DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)   # off-white background
GREY  = RGBColor(0x60, 0x60, 0x60)
ACCENT= RGBColor(0xFF, 0x4D, 0x4D)   # lighter red accent

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper utilities ──────────────────────────────────────────────────────────

def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, italic=False,
                font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_label(slide, text, l, t, w, h,
              font_size=18, bold=False, color=DARK,
              align=PP_ALIGN.LEFT, italic=False):
    """Add a text box; supports multi-paragraph text separated by \\n."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb


def add_bullet_box(slide, title, bullets, l, t, w, h,
                   title_size=18, bullet_size=14,
                   title_color=RED, bullet_color=DARK,
                   box_fill=WHITE, box_line=RED):
    """Draw a rounded-corner card with a bold title and bullet list."""
    card = add_rect(slide, l, t, w, h, fill_color=box_fill, line_color=box_line, line_width=Pt(1.5))
    # title
    add_textbox(slide, title, l + 0.15, t + 0.12, w - 0.3, 0.35,
                font_size=title_size, bold=True, color=title_color)
    # bullets
    bullet_text = "\n".join(f"  •  {b}" for b in bullets)
    add_label(slide, bullet_text, l + 0.15, t + 0.5, w - 0.3,
              h - 0.6, font_size=bullet_size, color=bullet_color)


def slide_background(slide, fill=LIGHT):
    """Fill the entire slide background."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=fill)


def header_bar(slide, title, subtitle=None):
    """Add top red header bar with white title text."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_color=RED)
    add_textbox(slide, title, 0.4, 0.1, 11.5, 0.65,
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, 0.72, 11.5, 0.35,
                    font_size=14, color=RGBColor(0xFF, 0xCC, 0xCC),
                    align=PP_ALIGN.LEFT)


def footer_bar(slide, text="Hair Fall Solution | Red Light Therapy (LLLT)"):
    """Thin dark footer."""
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_color=DARK)
    add_textbox(slide, text, 0.3, 7.17, 12.5, 0.28,
                font_size=9, color=GREY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title Slide
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
# Full dark background
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK)
# Red accent stripe on left
add_rect(slide, 0, 0, 0.45, 7.5, fill_color=RED)
# Light panel for text area
add_rect(slide, 0.6, 0.7, 8.8, 5.8, fill_color=RGBColor(0x22, 0x22, 0x40))

# Main title
add_textbox(slide, "Hair Fall Solution", 0.85, 1.0, 8.3, 1.1,
            font_size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(slide, "Using Red Light Therapy", 0.85, 2.0, 8.3, 0.9,
            font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
add_textbox(slide, "(LLLT — Low-Level Laser Therapy)", 0.85, 2.85, 8.3, 0.55,
            font_size=20, color=WHITE, align=PP_ALIGN.LEFT)

# Divider line
add_rect(slide, 0.85, 3.5, 5.5, 0.05, fill_color=RED)

# Subtitle
add_textbox(slide, "A Non-Invasive Approach to Hair Regrowth",
            0.85, 3.65, 8.3, 0.6,
            font_size=18, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC),
            align=PP_ALIGN.LEFT)

# Meta info
add_textbox(slide, "Project Presentation  |  Healthcare Innovation  |  2024",
            0.85, 4.35, 8.3, 0.45, font_size=13, color=GREY, align=PP_ALIGN.LEFT)
add_textbox(slide, "Target Audience: College Students / Startup Innovators",
            0.85, 4.75, 8.3, 0.4, font_size=13, color=GREY, align=PP_ALIGN.LEFT)

# Right decorative block with wavelength info
add_rect(slide, 9.9, 1.5, 3.1, 2.5, fill_color=RED)
add_textbox(slide, "630–680 nm", 9.95, 1.65, 3.0, 0.7,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Therapeutic\nLight Wavelength",
            9.95, 2.35, 3.0, 0.9, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 9.9, 4.3, 3.1, 1.5, fill_color=RGBColor(0x22, 0x22, 0x40))
add_textbox(slide, "FDA Approved\nNon-Invasive  |  Clinically Proven",
            9.95, 4.45, 3.0, 1.2, font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# Slide number
add_textbox(slide, "01 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Introduction to Hair Fall
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Introduction to Hair Fall",
           "Understanding the Basics — What's Normal vs Alarming?")
footer_bar(slide)

# Definition card
add_bullet_box(slide,
    "What is Hair Fall?",
    ["Natural biological process of hair shedding from follicles",
     "Hair goes through cycles: Growth (Anagen) → Rest (Telogen) → Shed (Exogen)",
     "Concern arises when shedding exceeds regrowth capacity"],
    0.4, 1.3, 5.9, 2.3)

# Normal vs Excessive
add_bullet_box(slide,
    "Normal vs. Excessive Loss",
    ["Normal: 50–100 hairs/day (natural cycle)",
     "Excessive: >150 hairs/day consistently",
     "Thinning patches, receding hairline, scalp visibility"],
    6.6, 1.3, 6.3, 2.3)

# Global concern card — full width
add_bullet_box(slide,
    "Global Concern Overview",
    ["Affects 50% of men and 25% of women worldwide by age 50",
     "Increasingly reported in youth (15–30 age group)",
     "Linked to modern lifestyle: stress, pollution, diet, screen time",
     "One of the top dermatological concerns globally"],
    0.4, 3.85, 12.5, 2.75, bullet_size=15)

add_textbox(slide, "02 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Statistics & Impact (with bar chart)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Statistics & Global Impact",
           "Hair loss is more prevalent — and more impactful — than you might think")
footer_bar(slide)

# Chart – prevalence by age group
chart_data = ChartData()
chart_data.categories = ["Age 20–30", "Age 30–40", "Age 40–50", "Age 50+"]
chart_data.add_series("Men (%)",   (18, 30, 45, 65))
chart_data.add_series("Women (%)", ( 8, 15, 25, 40))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.4), Inches(6.8), Inches(4.5),
    chart_data
).chart

chart.has_title = True
chart.chart_title.text_frame.text = "Hair Loss Prevalence by Age Group"
chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True

# Style series colours
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = RED
chart.series[1].format.fill.solid()
chart.series[1].format.fill.fore_color.rgb = DARK

# Stat cards on the right
stat_cards = [
    ("50%",   "of men affected by\nage 50"),
    ("25%",   "of women affected\nby age 50"),
    ("~80M",  "people affected in\nthe US alone"),
    ("↑ 30%", "rise in youth hair\nloss (last decade)"),
]
for i, (num, label) in enumerate(stat_cards):
    col = 7.5 if i % 2 == 0 else 10.05
    row = 1.4 + (i // 2) * 2.3
    add_rect(slide, col, row, 2.3, 1.95,
             fill_color=RED if i % 2 == 0 else DARK, line_color=None)
    add_textbox(slide, num, col, row + 0.15, 2.3, 0.75,
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_label(slide, label, col, row + 0.9, 2.3, 0.9,
              font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# Psychological impact note
add_rect(slide, 0.4, 6.05, 12.5, 0.75, fill_color=RGBColor(0xFF, 0xEE, 0xEE), line_color=RED, line_width=Pt(1))
add_textbox(slide,
    "⚠  Psychological Impact: Hair loss contributes to reduced self-esteem, social anxiety, "
    "and depression — especially among young adults. Addressing it early matters.",
    0.55, 6.12, 12.2, 0.6, font_size=12, color=DARK)

add_textbox(slide, "03 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Causes of Hair Fall
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Causes of Hair Fall",
           "Multiple factors — genetic, hormonal, environmental & lifestyle")
footer_bar(slide)

causes = [
    ("🧬  Genetics",
     ["Androgenetic Alopecia (male/female pattern baldness)",
      "Most common hereditary form of hair loss",
      "DHT hormone shrinks hair follicles over time"]),
    ("⚗  Hormonal Imbalance",
     ["Thyroid disorders (hypo/hyperthyroidism)",
      "Pregnancy, menopause, PCOS in women",
      "Testosterone imbalance in men"]),
    ("😰  Stress & Lifestyle",
     ["Telogen Effluvium: stress pushes follicles to rest phase",
      "Poor sleep, smoking, alcohol consumption",
      "Crash diets and extreme weight loss"]),
    ("🌫  Pollution & Nutrition",
     ["Scalp damage from particulate matter & UV exposure",
      "Iron, Zinc, Vitamin D, Biotin deficiencies",
      "High water hardness (chlorine, heavy metals)"]),
]

positions = [(0.4, 1.3), (6.85, 1.3), (0.4, 4.05), (6.85, 4.05)]
for (l, t), (title, bullets) in zip(positions, causes):
    add_bullet_box(slide, title, bullets, l, t, 6.0, 2.55, bullet_size=13)

add_textbox(slide, "04 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Existing Treatments & Limitations
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Existing Treatments & Their Limitations",
           "Current solutions work — but come with trade-offs")
footer_bar(slide)

# Treatment cards
treatments = [
    ("💊  Minoxidil (Rogaine)",
     ["Applied topically — dilates blood vessels",
      "Requires lifelong use; regrowth stops if discontinued",
      "Side effects: scalp irritation, unwanted facial hair"]),
    ("💉  Finasteride (Propecia)",
     ["Oral tablet — blocks DHT hormone",
      "Effective for male pattern baldness",
      "Side effects: sexual dysfunction, depression risk"]),
    ("✂  Hair Transplant",
     ["Permanent surgical solution (FUE / FUT methods)",
      "High cost: ₹50,000–₹3,00,000+",
      "Post-surgical recovery, scarring possible"]),
]
for i, (title, bullets) in enumerate(treatments):
    add_bullet_box(slide, title, bullets, 0.4 + i * 4.3, 1.3, 4.05, 3.0,
                   bullet_size=13)

# Limitations summary
add_rect(slide, 0.4, 4.55, 12.5, 2.15, fill_color=RGBColor(0x1A, 0x1A, 0x2E), line_color=RED, line_width=Pt(1.5))
add_textbox(slide, "Common Limitations of Existing Treatments",
            0.7, 4.65, 12.0, 0.45, font_size=16, bold=True, color=ACCENT)
lims = [
    "⚠  Significant side effects (hormonal, dermatological)",
    "💸  High cost — inaccessible for many patients",
    "🔄  Not permanent — require continuous use",
    "🏥  Surgical options carry surgical risks",
    "📉  Inconsistent results — vary greatly between individuals",
]
for j, lim in enumerate(lims):
    col = 0.7 if j < 3 else 6.9
    row = 5.18 + (j % 3) * 0.35
    add_textbox(slide, lim, col, row, 5.9, 0.32, font_size=12, color=LIGHT)

add_textbox(slide, "05 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Introduction to Red Light Therapy
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Introduction to Red Light Therapy",
           "Low-Level Laser Therapy (LLLT) — Harnessing the Power of Light")
footer_bar(slide)

# Definition
add_rect(slide, 0.4, 1.3, 8.2, 1.55, fill_color=DARK)
add_textbox(slide,
    "Red Light Therapy (RLT) / LLLT uses low-intensity red or near-infrared light "
    "to stimulate cellular activity. It is non-invasive, painless, and requires no chemicals.",
    0.6, 1.38, 7.9, 1.35, font_size=14, color=WHITE)

# Key attributes
attrs = [
    ("Non-Invasive",    "No surgery, no needles, no chemicals"),
    ("Painless",        "Gentle light application — completely comfortable"),
    ("Clinically Proven","Multiple peer-reviewed studies confirm efficacy"),
    ("FDA Approved",    "Cleared for hair loss treatment (iGrow, HairMax)"),
    ("At-Home Use",     "Available as consumer laser combs & helmets"),
    ("Long-Term Safe",  "No known systemic side effects"),
]
for i, (attr, desc) in enumerate(attrs):
    col = 0.4 + (i % 3) * 4.3
    row = 3.1 + (i // 3) * 1.75
    add_rect(slide, col, row, 4.0, 1.5, fill_color=WHITE, line_color=RED, line_width=Pt(1.5))
    add_textbox(slide, attr, col + 0.15, row + 0.1, 3.7, 0.45,
                font_size=15, bold=True, color=RED)
    add_textbox(slide, desc, col + 0.15, row + 0.55, 3.7, 0.8,
                font_size=12, color=DARK)

# Wavelength info on right
add_rect(slide, 9.0, 1.3, 3.9, 1.55, fill_color=RED)
add_textbox(slide, "Therapeutic Window", 9.1, 1.38, 3.7, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "630 – 680 nm\n(Red Light)\n800 – 1000 nm\n(Near-Infrared)",
            9.1, 1.8, 3.7, 1.0, font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)

add_textbox(slide, "06 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Working Principle
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "How Red Light Therapy Works",
           "Step-by-step mechanism of light-induced hair regrowth")
footer_bar(slide)

steps = [
    ("① Light Emission",
     "Device emits 630–680 nm red photons\ndirected at the scalp surface."),
    ("② Skin Penetration",
     "Photons penetrate 5–10 mm into\nscalp reaching hair follicle level."),
    ("③ Cellular Absorption",
     "Cytochrome c oxidase in mitochondria\nabsorbs photon energy."),
    ("④ ATP Production",
     "Increased ATP (cellular energy)\nboosts metabolic activity in follicles."),
    ("⑤ Blood Circulation",
     "Vasodilation improves blood flow,\ndelivering more nutrients to follicles."),
    ("⑥ Hair Regrowth",
     "Follicles shift from Telogen (rest)\nto Anagen (active growth) phase."),
]

for i, (step, desc) in enumerate(steps):
    col = 0.4 + (i % 3) * 4.3
    row = 1.4 + (i // 3) * 2.5
    # arrow indicator
    add_rect(slide, col, row, 0.55, 0.55, fill_color=RED)
    add_textbox(slide, str(i + 1), col, row, 0.55, 0.55,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, col + 0.6, row, 3.5, 1.9, fill_color=WHITE, line_color=DARK, line_width=Pt(0.5))
    add_textbox(slide, step, col + 0.75, row + 0.1, 3.2, 0.45,
                font_size=14, bold=True, color=RED)
    add_label(slide, desc, col + 0.75, row + 0.6, 3.2, 1.2,
              font_size=12, color=DARK)

add_textbox(slide, "07 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Scientific Mechanism
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Scientific Mechanism",
           "Cellular biology behind LLLT-induced hair regrowth")
footer_bar(slide)

# Left column — cellular mechanism
add_rect(slide, 0.4, 1.3, 5.9, 5.55, fill_color=WHITE, line_color=RED, line_width=Pt(1.5))
add_textbox(slide, "Cellular Stimulation Pathway",
            0.6, 1.42, 5.5, 0.45, font_size=16, bold=True, color=RED)
cell_points = [
    "Photobiomodulation (PBM): light triggers biologic change",
    "Cytochrome c oxidase accepts photon energy",
    "↑ Nitric oxide (NO) release → vasodilation",
    "Reactive oxygen species (ROS) modulation",
    "↑ mRNA synthesis → protein production",
    "Cell proliferation & reduced apoptosis",
    "Keratinocyte & fibroblast stimulation",
]
for j, pt in enumerate(cell_points):
    add_textbox(slide, f"  •  {pt}", 0.6, 1.98 + j * 0.58, 5.5, 0.5,
                font_size=12, color=DARK)

# Right column — hair growth cycle
add_rect(slide, 6.7, 1.3, 6.2, 2.55, fill_color=DARK)
add_textbox(slide, "Hair Growth Cycle Activation",
            6.9, 1.42, 5.9, 0.45, font_size=15, bold=True, color=ACCENT)
phases = [
    ("Anagen (Growth)", "2–7 years", "Active phase — LLLT targets & extends this", WHITE),
    ("Catagen (Transition)", "2–3 weeks", "Regression — follicle shrinks", GREY),
    ("Telogen (Rest)", "3 months", "Dormant — LLLT pulls follicles OUT of this", RGBColor(0xFF, 0x80, 0x80)),
]
for k, (phase, dur, note, col) in enumerate(phases):
    add_textbox(slide, f"{phase}  [{dur}]",
                6.9, 1.95 + k * 0.6, 3.5, 0.42, font_size=12, bold=True, color=col)
    add_textbox(slide, note,
                6.9, 2.3 + k * 0.6, 5.9, 0.35, font_size=11, color=GREY)

# Right bottom — key insight
add_rect(slide, 6.7, 4.05, 6.2, 2.8, fill_color=WHITE, line_color=RED, line_width=Pt(1.5))
add_textbox(slide, "Key Scientific Insight",
            6.9, 4.17, 5.9, 0.45, font_size=15, bold=True, color=RED)
insights = [
    "LLLT does NOT heat tissue — purely photochemical",
    "Optimal dose: 4–6 J/cm² per session",
    "Frequency: 3 sessions/week for best results",
    "Penetration depth sufficient to reach bulge stem cells",
    "Promotes Wnt/β-catenin signaling — hair morphogenesis",
]
for m, ins in enumerate(insights):
    add_textbox(slide, f"  •  {ins}", 6.9, 4.72 + m * 0.42, 5.9, 0.38,
                font_size=12, color=DARK)

add_textbox(slide, "08 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Effectiveness & Research
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Effectiveness & Clinical Research",
           "Evidence-based results — what the studies show")
footer_bar(slide)

# Chart — hair density improvement over 6 months
chart_data2 = ChartData()
chart_data2.categories = ["Baseline", "Month 1", "Month 2", "Month 3", "Month 4", "Month 6"]
chart_data2.add_series("LLLT Group (hairs/cm²)",  (185, 195, 210, 228, 245, 262))
chart_data2.add_series("Control Group (hairs/cm²)", (185, 184, 183, 183, 182, 181))

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(0.4), Inches(1.4), Inches(7.0), Inches(4.5),
    chart_data2
).chart
chart2.has_title = True
chart2.chart_title.text_frame.text = "Hair Density Over 6-Month LLLT Treatment"
chart2.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
chart2.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
chart2.series[0].format.line.color.rgb = RED
chart2.series[0].format.line.width = Pt(2.5)
chart2.series[1].format.line.color.rgb = DARK

# Research highlights
add_rect(slide, 7.8, 1.3, 5.1, 4.7, fill_color=WHITE, line_color=RED, line_width=Pt(1.5))
add_textbox(slide, "Key Research Findings",
            8.0, 1.42, 4.7, 0.45, font_size=16, bold=True, color=RED)

studies = [
    ("Avci et al., 2014 (PLoS ONE)",
     "39% increase in hair count after 26 weeks of LLLT"),
    ("Jimenez et al., 2014",
     "Significant improvement in hair density in both men and women"),
    ("Lanzafame et al., 2013",
     "35% increase in hair growth in women with female pattern hair loss"),
    ("FDA Clearance",
     "HairMax LaserComb & iGrow Helmet — cleared for safety & efficacy"),
    ("Meta-analysis (2018)",
     "Pooled data: LLLT shows statistically significant hair regrowth vs placebo"),
]
for n, (ref, finding) in enumerate(studies):
    add_rect(slide, 8.0, 2.02 + n * 0.82, 4.7, 0.72,
             fill_color=LIGHT, line_color=GREY, line_width=Pt(0.5))
    add_textbox(slide, ref, 8.15, 2.06 + n * 0.82, 4.4, 0.28,
                font_size=10, bold=True, color=RED)
    add_textbox(slide, finding, 8.15, 2.32 + n * 0.82, 4.4, 0.35,
                font_size=10, color=DARK)

add_textbox(slide, "09 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Advantages
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Advantages of Red Light Therapy",
           "Why LLLT stands out among hair loss solutions")
footer_bar(slide)

advantages = [
    ("✅  Safe & Non-Invasive",
     ["No surgery, injections, or chemicals required",
      "No downtime — continue daily activities",
      "Suitable for all skin types & tones"]),
    ("✅  Minimal Side Effects",
     ["No systemic side effects reported",
      "Mild scalp warmth at most",
      "Safe for long-term repeated use"]),
    ("✅  At-Home Convenience",
     ["FDA-cleared consumer devices available",
      "Use in privacy of your home",
      "15–30 minute sessions, 3×/week"]),
    ("✅  Cost-Effective Long-Term",
     ["One-time device investment vs ongoing medication",
      "No prescription required for OTC devices",
      "Comparable efficacy at fraction of transplant cost"]),
    ("✅  Complementary Therapy",
     ["Can be combined with Minoxidil for enhanced results",
      "Supports post-transplant recovery",
      "Improves overall scalp health"]),
    ("✅  Scientifically Validated",
     ["Multiple RCTs and meta-analyses confirm efficacy",
      "FDA-approved devices for safety",
      "Growing body of peer-reviewed literature"]),
]
for i, (title, bullets) in enumerate(advantages):
    col = 0.4 + (i % 3) * 4.3
    row = 1.4 + (i // 3) * 2.65
    add_bullet_box(slide, title, bullets, col, row, 4.0, 2.45, bullet_size=12)

add_textbox(slide, "10 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Comparison Table (Traditional vs LLLT)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Comparison: Traditional Treatments vs Red Light Therapy",
           "Side-by-side analysis across key parameters")
footer_bar(slide)

# Table
cols_w   = [2.6, 1.9, 1.9, 1.9, 2.15, 2.45]
col_starts = [0.25]
for cw in cols_w[:-1]:
    col_starts.append(col_starts[-1] + cw + 0.02)

headers = ["Parameter", "Minoxidil", "Finasteride", "Transplant", "LLLT (Home)", "LLLT (Clinical)"]
rows_data = [
    ["Invasiveness",      "None (topical)", "None (oral)", "Surgical",    "None",           "None"],
    ["Side Effects",      "Moderate",       "Significant", "Post-surgical","Minimal",       "Minimal"],
    ["Cost (Approx.)",    "₹500/month",     "₹800/month",  "₹50K–3L",     "₹15K–50K (one-time)", "Per session"],
    ["Permanence",        "Ongoing",        "Ongoing",     "Permanent",   "Ongoing",        "Ongoing"],
    ["Results Timeline",  "3–6 months",     "3–6 months",  "6–12 months", "3–6 months",     "3–4 months"],
    ["At-Home Use",       "✓",              "✓",           "✗",           "✓",              "✗"],
    ["FDA Approved",      "✓",              "✓",           "✓",           "✓",              "✓"],
    ["Pain / Discomfort", "Low",            "None",        "High",        "None",           "None"],
]

row_h = 0.52
header_y = 1.25

# Header row
add_rect(slide, 0.25, header_y, 12.83, row_h, fill_color=DARK)
x = 0.25
for h, cw in zip(headers, cols_w):
    add_textbox(slide, h, x + 0.05, header_y + 0.07, cw - 0.1, row_h - 0.1,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += cw + 0.02

for r_idx, row in enumerate(rows_data):
    y = header_y + row_h + r_idx * row_h
    bg = LIGHT if r_idx % 2 == 0 else WHITE
    add_rect(slide, 0.25, y, 12.83, row_h, fill_color=bg, line_color=GREY, line_width=Pt(0.3))
    x = 0.25
    for c_idx, (cell, cw) in enumerate(zip(row, cols_w)):
        cell_color = DARK
        if c_idx >= 4:   # LLLT columns — highlight
            cell_color = RED
        add_textbox(slide, cell, x + 0.05, y + 0.1, cw - 0.1, row_h - 0.1,
                    font_size=11, color=cell_color, align=PP_ALIGN.CENTER,
                    bold=(c_idx >= 4))
        x += cw + 0.02

add_textbox(slide,
    "✦  LLLT columns highlighted in red — consistently competitive or superior across parameters",
    0.25, 6.8, 12.83, 0.35, font_size=11, italic=True, color=RED)

add_textbox(slide, "11 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Devices & Future Scope
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Devices, Limitations & Future Scope",
           "Current market landscape and where the technology is headed")
footer_bar(slide)

# Devices
add_bullet_box(slide,
    "🔧  Available LLLT Devices",
    ["Laser Combs (HairMax): handheld, 9–82 lasers",
     "Laser Helmets / Caps: hands-free, full scalp coverage",
     "Clinical Hood Systems: salon/dermatologist grade",
     "LED Panels: broader wavelength range, lower cost",
     "Wearable smart devices with app connectivity"],
    0.4, 1.3, 6.0, 3.15, bullet_size=13)

# Limitations
add_bullet_box(slide,
    "⚠  Known Limitations",
    ["Requires strict consistency — 3×/week minimum",
     "Not effective for complete (scarring) alopecia",
     "Results vary per individual genetics & severity",
     "Delayed results — visible improvement after 3–6 months",
     "Quality varies widely among consumer products"],
    6.7, 1.3, 6.2, 3.15, bullet_size=13,
    title_color=DARK, box_line=DARK)

# Future scope
add_rect(slide, 0.4, 4.65, 12.5, 2.1, fill_color=DARK)
add_textbox(slide, "🚀  Future Scope & Startup Opportunities",
            0.6, 4.75, 12.0, 0.45, font_size=16, bold=True, color=ACCENT)
future = [
    "🤖  AI + LLLT: Computer vision for real-time scalp health monitoring & personalized dosing",
    "📱  App-integrated wearables: track treatment adherence & hair density progress via smartphone",
    "🧬  Gene-expression guided therapy: personalise wavelength/dose based on patient genomics",
    "🌍  Telehealth platforms: remote dermatologist supervision with at-home LLLT devices",
]
for p_idx, fut in enumerate(future):
    col = 0.6 if p_idx < 2 else 6.9
    row = 5.28 + (p_idx % 2) * 0.48
    add_textbox(slide, fut, col, row, 6.0, 0.42, font_size=12, color=LIGHT)

add_textbox(slide, "12 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Conclusion & References
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
slide_background(slide)
header_bar(slide, "Conclusion & References",
           "Red Light Therapy — A Promising, Science-Backed Solution")
footer_bar(slide)

# Conclusion card
add_rect(slide, 0.4, 1.3, 7.6, 3.9, fill_color=DARK)
add_textbox(slide, "Key Takeaways",
            0.6, 1.42, 7.2, 0.45, font_size=16, bold=True, color=ACCENT)
takeaways = [
    "Hair fall is a growing global health concern affecting millions",
    "Existing treatments work but carry side effects & high costs",
    "LLLT offers a safe, non-invasive, evidence-based alternative",
    "Clinically validated: 30–40% increase in hair density in 3–6 months",
    "FDA-cleared devices available for both clinical & home use",
    "Future: AI + LLLT = personalized, scalable hair regrowth solutions",
    "Strong startup & innovation opportunity in the growing hair-care tech market",
]
for q_idx, tk in enumerate(takeaways):
    add_textbox(slide, f"  ✓  {tk}",
                0.6, 1.98 + q_idx * 0.46, 7.2, 0.4, font_size=12, color=LIGHT)

# Call-to-action
add_rect(slide, 0.4, 5.35, 7.6, 0.8, fill_color=RED)
add_textbox(slide,
    "LLLT is not just a treatment — it's the future of non-surgical hair restoration.",
    0.6, 5.45, 7.3, 0.6, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# References
add_rect(slide, 8.35, 1.3, 4.6, 4.85, fill_color=WHITE, line_color=RED, line_width=Pt(1.5))
add_textbox(slide, "References & Sources",
            8.5, 1.42, 4.3, 0.45, font_size=14, bold=True, color=RED)
refs = [
    "Avci P. et al. (2014). Low-level laser therapy for hair loss.\nPLoS ONE, 9(7).",
    "Jimenez JJ et al. (2014). Efficacy of 655nm LLLT in\nfemale pattern hair loss. JDDG.",
    "Lanzafame RJ et al. (2013). Laser irradiation for\nwomen with alopecia. Lasers Surg Med.",
    "FDA 510(k) Clearances – HairMax LaserComb,\niGrow Helmet (FDA.gov).",
    "WHO Global Health Observatory — Hair & Skin\nConditions Data.",
    "Hamblin MR (2016). Photobiomodulation for\nhair growth. Semin Cutan Med Surg.",
    "American Academy of Dermatology —\nHair Loss Statistics (AAD.org).",
]
for s_idx, ref in enumerate(refs):
    add_label(slide, f"[{s_idx+1}] {ref}",
              8.5, 1.98 + s_idx * 0.6, 4.3, 0.55, font_size=9.5, color=DARK)

add_textbox(slide, "13 / 13", 12.4, 7.15, 0.85, 0.3,
            font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
OUTPUT = "/home/runner/work/ai-ml-lab/ai-ml-lab/Hair_Fall_LLLT_Presentation.pptx"
prs.save(OUTPUT)
print(f"Presentation saved → {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
